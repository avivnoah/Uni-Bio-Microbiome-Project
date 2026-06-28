"""Build deliverables/milestone1.pptx — paper/ink/rose "bold mono / editorial"
deck matching milestone1.html (unified: our framing/design + partner EDA + dual baseline).

Fonts: Space Grotesk (display) + JetBrains Mono (labels). PowerPoint substitutes
if the viewer lacks them; install from fontsource for an exact match.
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "analysis" / "figures"
_ALL = {}
for _p in sorted((ROOT / "analysis" / "findings").glob("*.json")):
    _ALL.update(json.loads(_p.read_text()))
F = PF = _ALL

PAPER = RGBColor(0xFA, 0xF9, 0xF6)
INK = RGBColor(0x0A, 0x0A, 0x0A)
INK2 = RGBColor(0x1C, 0x1B, 0x19)
MUT = RGBColor(0x6B, 0x68, 0x63)
LINE = RGBColor(0xDA, 0xD7, 0xCF)
ACCENT = RGBColor(0xFF, 0x2D, 0x55)
ROSE_GHOST = RGBColor(0xF1, 0xC9, 0xD2)
DISP = "Space Grotesk"
MONO = "JetBrains Mono"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height
ML, MR = 0.85, 0.85
TOTAL = 16


def slide_bg(s, color=PAPER):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, runs, size, font=DISP, color=INK, bold=False, first=False,
         align=PP_ALIGN.LEFT, spacing=1.0, space_after=6, tracking=None, caps=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        if caps:
            text = text.upper()
        r = p.add_run()
        r.text = text
        r.font.name = ov.get("font", font)
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.italic = ov.get("italic", False)
        r.font.color.rgb = ov.get("color", color)
        if tracking is not None:
            _track(r, tracking)
    return p


def _track(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def line(s, x, y, w, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def vline(s, x, y, h, color=LINE, weight=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def rect(s, x, y, w, h, color=ACCENT):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def shell(s, section, n, num_pos="br"):
    slide_bg(s)
    if num_pos == "br":
        nt = box(s, 8.7, 4.0, 5.6, 3.9)
        para(nt, f"{n:02d}", 250, font=DISP, color=ROSE_GHOST, bold=True, first=True, align=PP_ALIGN.RIGHT)
    rect(s, ML, 0.46, 0.10, 0.10, ACCENT)
    ht = box(s, ML + 0.22, 0.4, 7.0, 0.3)
    para(ht, "METACARDIS · MILESTONE 1", 9.5, font=MONO, color=MUT, first=True, tracking=1.6)
    hr = box(s, W.inches - MR - 6.0, 0.4, 6.0, 0.3)
    para(hr, section.upper(), 9.5, font=MONO, color=MUT, first=True, align=PP_ALIGN.RIGHT, tracking=1.6)
    line(s, ML, 0.82, W.inches - ML - MR, LINE, 1.0)
    line(s, ML, 6.86, W.inches - ML - MR, LINE, 1.0)
    ft = box(s, ML, 6.95, 7.0, 0.3)
    para(ft, "COMPUTATIONAL MICROBIOME WORKSHOP", 9, font=MONO, color=MUT, first=True, tracking=1.4)
    fr = box(s, W.inches - MR - 3.0, 6.95, 3.0, 0.3)
    para(fr, f"{n:02d} / {TOTAL:02d}", 9, font=MONO, color=MUT, first=True, align=PP_ALIGN.RIGHT, tracking=1.4)


def eyebrow(s, x, y, text, w=11.0):
    tf = box(s, x, y, w, 0.35)
    para(tf, text, 11, font=MONO, color=ACCENT, bold=False, first=True, tracking=1.8)


def new():
    return prs.slides.add_slide(BLANK)


def pic(s, name, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(str(FIG / name), Inches(x), Inches(y), **kw)


def stat(s, x, y, num_runs, label, w=2.9, big=44):
    tf = box(s, x, y, w, 1.0)
    para(tf, num_runs, big, font=DISP, color=INK, bold=True, first=True, spacing=0.95)
    lf = box(s, x, y + 0.78, w, 0.5)
    para(lf, label, 9, font=MONO, color=MUT, first=True, tracking=1.2)


def mini_table(s, x, y, headers, rows, colw, valsize=16):
    for j, (lab, al) in enumerate(headers):
        hb = box(s, x + sum(colw[:j]), y, colw[j], 0.3)
        para(hb, lab, 9, font=MONO, color=MUT, first=True, align=al, tracking=1.0)
    tw = sum(colw)
    line(s, x, y + 0.34, tw, INK, 1.6)
    ry = y + 0.5
    for lab, vals in rows:
        lb = box(s, x, ry, colw[0], 0.45)
        para(lb, lab, 11, font=MONO, color=MUT, first=True)
        for j, (v, win) in enumerate(vals):
            vb = box(s, x + sum(colw[:j + 1]), ry - 0.03, colw[j + 1], 0.45)
            para(vb, f"{v}", valsize, font=MONO, color=(ACCENT if win else INK),
                 bold=win, first=True, align=PP_ALIGN.RIGHT)
        line(s, x, ry + 0.5, tw, LINE, 1.0)
        ry += 0.62
    return ry


ms, gc = F["modality_structure"], F["group_counts"]
pcoa, pg, pc = F["pcoa_explained"], F["permanova_group"], F["permanova_center"]
dp, b = F["disease_pred_balacc"], F["naive_distance_baseline"]
fi, md_acc = PF["metadata_feature_importance"], PF["metadata_only_cv_balacc"]
dww = PF["disease_with_without"]
pnaive = PF["partner_naive_imputation_r2"]
max_sp = max(PF["diff_species_per_group"].values())
max_mb = max(PF["diff_metab_per_group"].values())


# ---- 01 title ----
s = new(); slide_bg(s)
nt = box(s, -0.2, -0.4, 6.0, 3.5)
para(nt, "01", 250, font=DISP, color=ROSE_GHOST, bold=True, first=True)
rect(s, ML, 0.46, 0.10, 0.10, ACCENT)
ht = box(s, ML + 0.22, 0.4, 7.0, 0.3)
para(ht, "METACARDIS · MILESTONE 1", 9.5, font=MONO, color=MUT, first=True, tracking=1.6)
line(s, ML, 0.82, W.inches - ML - MR, LINE, 1.0)
eyebrow(s, ML, 1.7, "MULTI-OMICS · MISSING-MODALITY DISTANCES")
tf = box(s, ML, 2.15, 9.8, 3.2)
para(tf, "Estimating sample distances", 40, font=DISP, color=INK, bold=True, first=True, spacing=1.0)
para(tf, "across multi-omics,", 40, font=DISP, color=INK, bold=True, spacing=1.0)
para(tf, [("under missing modalities", {"color": ACCENT})], 40, font=DISP, bold=True, spacing=1.0)
rect(s, ML, 5.25, 1.15, 0.05, ACCENT)
bf = box(s, ML, 5.5, 9.0, 0.4)
para(bf, "Group ____________   ·   MetaCardis cohort", 11, font=MONO, color=MUT, first=True, tracking=0.8)

# ---- 02 statement ----
s = new(); shell(s, "The problem", 2)
eyebrow(s, ML, 1.6, "WHAT WE ARE ACTUALLY ASKED TO BUILD")
tf = box(s, ML, 2.1, 10.6, 1.7)
para(tf, [("A sample×sample ", {}), ("distance matrix", {"color": ACCENT, "bold": True}),
          (" from two omics — that holds up when an entire modality is ", {}),
          ("missing", {"color": ACCENT, "bold": True}), (".", {})],
     30, font=DISP, color=INK, first=True, spacing=1.04)
bf = box(s, ML, 4.4, 11.2, 2.0)
para(bf, [("Graded by ", {}), ("Mantel test", {"bold": True}),
          (" vs. ground truth: standardise all features → PCA → Euclidean in the first 2 PCs. The target geometry is linear.", {})],
     15, font=DISP, color=INK2, first=True, spacing=1.3, space_after=10)
para(bf, [("With both modalities present that target is trivially recomputable — so the whole problem lives in the ", {}),
          ("missing-modality regime", {"color": ACCENT, "bold": True}), (".", {})],
     15, font=DISP, color=INK2, spacing=1.3)

# ---- 03 data overview ----
s = new(); shell(s, "Exploration / data", 3, num_pos="none")
eyebrow(s, ML, 1.1, "THE DATA ALREADY CONTAINS THE CHALLENGE")
tf = box(s, ML, 1.5, 11.0, 0.8)
para(tf, "~40% of subjects are missing a whole modality", 26, font=DISP, color=INK, bold=True, first=True)
pic(s, "01_overview.png", ML, 2.45, w=11.63)
line(s, ML, 5.75, W.inches - ML - MR, LINE, 1.0)
xs = ML
for num, lab in [([("1042", {})], "BOTH OMICS"), ([("348", {})], "MICROBIOME ONLY"),
                 ([("348", {})], "METABOLOME ONLY"),
                 ([("528", {}), (" → ", {"color": ACCENT}), ("18", {})], "T2D → CHF (IMBALANCE)")]:
    stat(s, xs, 5.95, num, lab, w=2.9, big=40)
    xs += 2.95

# ---- 04 demographics confound ----
s = new(); shell(s, "Exploration / confounders", 4, num_pos="none")
eyebrow(s, ML, 1.15, "DEMOGRAPHICS DIFFER SHARPLY BETWEEN DISEASE GROUPS")
pic(s, "10_age_bmi_by_group.png", ML, 1.75, w=11.63)
fn = box(s, ML, 6.25, 11.6, 0.4)
para(fn, "OBESITY 2A/2B — YOUNGEST, BMI > 40 · DIABETES/CVD 3–6 — OLDER (~60) · AGE & BMI ARE PRIME CONFOUNDERS",
     9, font=MONO, color=MUT, first=True, tracking=0.6)

# ---- 05 batch ----
s = new(); shell(s, "Exploration / batch", 5)
eyebrow(s, ML, 1.15, "RECRUITMENT CENTER IS A CONFOUNDER, NOT NOISE")
pic(s, "12_country_by_group.png", ML, 1.9, w=6.7)
sx = 8.1
stat(s, sx, 2.0, [(f"{pc['F']}", {"color": ACCENT})], "PERMANOVA F · CENTER", w=4.4, big=46)
line(s, sx, 3.15, 3.6, LINE, 1.0)
stat(s, sx, 3.3, [(f"{pg['F']}", {"color": ACCENT})], "PERMANOVA F · DISEASE GROUP", w=4.4, big=46)
tf = box(s, sx, 4.55, 4.5, 2.0)
para(tf, [("Groups are country-skewed (2b≈German, 7≈French, 4 no-Germany) and ", {}),
          ("CENTER structures the microbiome more than disease", {"color": ACCENT, "bold": True}),
          (". Denmark is 100% missing demographics (MNAR).", {})],
     13, font=DISP, color=INK2, first=True, spacing=1.3)

# ---- 06 demographics out-predict omics ----
s = new(); shell(s, "Exploration / what predicts disease", 6)
eyebrow(s, ML, 1.15, "BASIC DEMOGRAPHICS PREDICT THE DISEASE BETTER THAN THE OMICS")
pic(s, "14_feature_importance.png", ML, 1.9, w=6.6)
sx = 8.0
stat(s, sx, 2.0, [(f"{md_acc}", {"color": ACCENT})], "METADATA-ONLY RF · 5-FOLD CV · BAL-ACC", w=4.5, big=52)
tf = box(s, sx, 3.4, 4.6, 2.9)
para(tf, [("Random Forest, 5-fold cross-validation, guessing the disease group (1 of 9). Same model, three inputs: metadata-only ", {}),
          (f"{md_acc}", {"bold": True}), (" · both-omics ", {}), (f"{dp['both_omics_complete']}", {"bold": True}),
          (" · microbiome-only ", {}), (f"{dp['microbiome_only_all']}", {"bold": True}),
          (" (chance≈0.11).", {})], 13.5, font=DISP, color=INK2, first=True, spacing=1.3, space_after=10)
para(tf, [("BMI (" + str(fi['BMI_C']) + ") and age (" + str(fi['AGE']) + ") dominate", {"color": ACCENT, "bold": True}),
          (" — the disease signal in the omics is weak and distributed, so distances (not labels) are the right target.", {})],
     13.5, font=DISP, color=INK2, spacing=1.3, space_after=10)
para(tf, [("Spec check — does imputing a missing omics help prediction? ", {"bold": True}),
          (f"No: microbiome alone {dww['without_imputation_microbiome_only']}, + an imputed metabolome "
           f"{dww['with_imputation_micro_plus_recovered_metab']}, + the real metabolome "
           f"{dww['real_both_omics_upper_bound']}.", {})],
     12, font=DISP, color=MUT, spacing=1.25)

# ---- 07 microbiome geometry ----
s = new(); shell(s, "Exploration / microbiome map", 7)
eyebrow(s, ML, 1.15, "DISEASE GROUPS DON'T FORM CLEAN CLUSTERS ON THE MICROBIOME MAP")
pic(s, "04b_microbiome_pcoa_disease.png", ML, 1.85, h=4.7)
sx = 8.4
stat(s, sx, 1.95, [(f"{pcoa[0]+pcoa[1]:.0%}", {"color": ACCENT})], "OF PATIENT DIFFERENCES VISIBLE IN THIS FLAT MAP · REST HIDDEN", w=4.1, big=46)
cf = box(s, sx, 3.15, 4.1, 0.4)
para(cf, f"PCOA (AITCHISON) · PC1 {pcoa[0]:.0%} · PC2 {pcoa[1]:.0%}", 9, font=MONO, color=MUT, first=True, tracking=0.8)
tf = box(s, sx, 3.6, 4.1, 2.9)
para(tf, [("Disease colours are mixed everywhere — ", {}), ("no clean clusters", {"color": ACCENT, "bold": True}),
          (", so we place patients by distance, not predict a label. The 2 axes capture only "
           f"{pcoa[0]:.0%}/{pcoa[1]:.0%}; structure is high-dimensional.", {})],
     13.5, font=DISP, color=INK2, first=True, spacing=1.3, space_after=10)
para(tf, "Country also shapes this map; we judge that with the df-adjusted PERMANOVA on slide 05, not by eye.",
     12, font=DISP, color=MUT, spacing=1.3)

# ---- 08 differential per group ----
s = new(); shell(s, "Exploration / differential", 8, num_pos="none")
eyebrow(s, ML, 1.15, "BOTH OMICS CARRY DISEASE SIGNAL — AT DIFFERENT SCALES")
pic(s, "15_diff_per_group.png", ML, 1.75, w=11.63)
fn = box(s, ML, 6.25, 11.6, 0.4)
para(fn, f"VS. CONTROL, FDR<0.05 · UP TO {max_sp} SPECIES AND {max_mb} METABOLITES DIFFER · METABOLOME SHIFTS ARE BROAD",
     9, font=MONO, color=MUT, first=True, tracking=0.6)

# ---- 09 cross-omics ----
s = new(); shell(s, "Exploration / coupling", 9)
eyebrow(s, ML, 1.15, "CROSS-OMICS COUPLING IS WEAK — THIS CAPS IMPUTATION")
pic(s, "09_cross_omics_corr.png", ML, 1.85, h=4.7)
sx = 7.8
stat(s, sx, 2.1, [(f"{F['max_abs_taxa_metab_corr']}", {"color": ACCENT})], "MAX |SPEARMAN| TAXA ↔ METABOLITE", w=4.6, big=52)
tf = box(s, sx, 3.5, 4.6, 2.8)
para(tf, [("Both partners find it independently (ρ < 0.25 in the top pairs). Only a metabolite-specific subset is predictable from microbiome → cross-modal recovery has a ", {}),
          ("low ceiling", {"color": ACCENT, "bold": True}), (".", {})],
     14, font=DISP, color=INK2, first=True, spacing=1.3)

# ---- 10 naive baseline (drop a whole modality, recover, recompute distances) ----
s = new(); shell(s, "Naïve baseline", 10)
eyebrow(s, ML, 1.15, "MANDATORY BASELINE — DROP A WHOLE MODALITY, RECOVER, THEN MEASURE DISTANCES")
vline(s, 6.85, 1.7, 3.6, LINE, 1.0)
lh = box(s, ML, 1.7, 6.0, 0.35)
para(lh, "DROP WHOLE METABOLOME · RECOVER · MANTEL", 10, font=MONO, color=ACCENT, bold=True, first=True, tracking=1.2)
left_rows = []
for frac in (10, 30, 50):
    d = b[f"drop_metab_{frac}pct"]
    m, med, k, rf = d["mean"]["mantel_r"], d["median"]["mantel_r"], d["knn"]["mantel_r"], d["rf"]["mantel_r"]
    best = max(m, med, k, rf)
    left_rows.append((f"{frac}%", [(m, m == best), (med, med == best), (k, k == best), (rf, rf == best)]))
mini_table(s, ML, 2.2, [("DROPPED", PP_ALIGN.LEFT), ("MEAN", PP_ALIGN.RIGHT), ("MED", PP_ALIGN.RIGHT),
           ("KNN", PP_ALIGN.RIGHT), ("RF", PP_ALIGN.RIGHT)], left_rows, [1.25, 1.0, 1.0, 1.0, 1.0], valsize=14)
ln = box(s, ML, 4.35, 5.7, 1.5)
para(ln, [("Row = share whose whole metabolome we hid; column = how we rebuilt it; number = Mantel score "
           "(1 = identical to the true distances).", {})], 11.5, font=DISP, color=MUT, first=True, spacing=1.2)
tf = box(s, 7.2, 1.95, 5.2, 3.4)
para(tf, [("All four ", {}), ("tie", {"color": ACCENT, "bold": True}),
          (" — once a whole metabolome is gone, the only way to rebuild it is from the microbiome, and the "
           "two are weakly linked (slide 09). So the clever fills add almost nothing over a plain mean; even "
           "mean clears ~0.77 at 50% missing.", {})],
     13, font=DISP, color=INK2, first=True, spacing=1.3, space_after=10)
para(tf, [("The high score is misleading: at low missing-% most patients keep both omics, so most distances "
           "are still exact and dominate it (slide 09). It is not the discriminating number.", {})],
     12.5, font=DISP, color=MUT, spacing=1.3)
mline = box(s, ML, 5.95, 11.6, 0.55)
para(mline, [("How the distances are made: ", {"bold": True, "color": INK}),
     ("each matrix is n×n (one row/column per patient), built the grader's way (standardize both omics, PCA, "
      "Euclidean distance in the first 2 PCs). The Mantel test scores the rank correlation between our matrix "
      "and the true one, over every patient pair.", {})],
     9, font=DISP, color=INK2, first=True, spacing=1.18)
dl = box(s, ML, 6.5, 11.6, 0.4)
para(dl, "ONE CONCLUSION: NAÏVE IMPUTATION CAN'T BEAT A TRIVIAL MEAN — THE CROSS-MODAL CEILING IS THE WALL THE PIPELINE MUST BREAK.",
     8.5, font=MONO, color=MUT, first=True, tracking=0.6)

# ---- 11 method framing ----
s = new(); shell(s, "Method selection", 11)
eyebrow(s, ML, 1.5, "HOW WE PICKED THE TWO METHODS ON THE NEXT SLIDES")
tf = box(s, ML, 1.95, 11.0, 0.9)
para(tf, "We scanned many fields, then tried hard to break what we found.", 26, font=DISP, color=INK, bold=True, first=True)
nf = box(s, ML, 3.0, 11.4, 2.3)
for lead, rest in [("Test 1", " — does it respect how we are scored? Mantel ranks distances against a PCA map, so the danger is not nonlinearity but methods like UMAP/t-SNE that distort the global geometry. We kept global-preserving methods."),
                   ("Test 2", " — can it place a patient missing a whole omics? We kept methods that build each omics' own similarity and combine them, so a patient is placeable from the omics they have.")]:
    para(nf, [("— ", {"color": ACCENT}), (lead, {"bold": True}), (rest, {})],
         13.5, font=DISP, color=INK2, spacing=1.25, space_after=10, first=(lead == "Test 1"))
rect(s, ML, 5.5, 0.04, 1.0, ACCENT)
cf = box(s, ML + 0.25, 5.5, 11.0, 1.1)
para(cf, [("Then we ", {}), ("tried to break our own picks", {"color": ACCENT, "bold": True}),
          (" and feature the hardened version that fixes each weakness. Every method must beat the naive baseline in the missing-omics case, or we report it as a negative result.", {})],
     14, font=DISP, color=INK, first=True, spacing=1.3)


def method_slide(n, section, eyebrow_txt, idx, title, bullets):
    s = new(); shell(s, section, n)
    eyebrow(s, ML, 1.45, eyebrow_txt)
    tf = box(s, ML, 1.9, 11.4, 1.0)
    para(tf, [(idx + "  ", {"color": ACCENT, "bold": True}), (title, {})],
         30, font=DISP, color=INK, bold=True, first=True)
    nf = box(s, ML, 3.15, 11.4, 3.4)
    for i, runs in enumerate(bullets):
        para(nf, [("— ", {"color": ACCENT})] + runs, 14.5, font=DISP, color=INK2,
             spacing=1.3, space_after=11, first=(i == 0))


method_slide(12, "Inspiring method ①", "OPTIMAL TRANSPORT · MATH & VISION ORIGIN · SCOT+ 2025", "①",
    "Augmented / Fused Gromov–Wasserstein",
    [[("GW compares samples in ", {}), ("incomparable spaces", {"bold": True}), (" using only within-modality geometry — no shared features needed.", {})],
     [("Adaptation: ", {"bold": True}), ("couple a missing-modality sample to the complete samples' geometry, then barycentric-project onto the joint-PCA anchor → distances.", {})],
     [("Honest fix (red-team): use ", {}), ("Augmented GW", {"color": ACCENT, "bold": True}), (" — the feature term breaks GW's rotation/reflection ambiguity.", {})],
     [("Gate: ", {"bold": True}), ("must beat mean/PLS in the missing regime, else reported as a negative result.", {})]])

method_slide(13, "Inspiring method ②", "WAYFINDING · GPS-STYLE TRIANGULATION · OUR SYNTHESIS", "②",
    "CM-Harmonic-Aniso: triangulate, don't impute",
    [[("Pin the 1042 complete patients at their ", {}), ("true first-two-PC coordinates", {"bold": True}), (" (the exact thing we are graded on) as fixed landmarks. The only method anchored on the grader's own target.", {})],
     [("Place each of the 696 single-omics patients by ", {}), ("triangulating off those landmarks", {"color": ACCENT, "bold": True}), (": one harmonic solve on a fused graph where a missing omics is no edge, never an invented value.", {})],
     [("Honest per-axis split: ", {"bold": True}), ("the part carried by the omics you have beats every imputation baseline; the part needing the missing omics falls back to the average, provably never worse than mean-impute.", {})],
     [("Gate: ", {"bold": True}), ("must beat mean/PLS per-axis on held-out patients, with a fail-closed fallback to PLS on boundary-poor cohorts, else a negative result.", {})]])

# ---- 14 pipeline + eval ----
s = new(); shell(s, "Pipeline & evaluation", 14)
eyebrow(s, ML, 1.3, "THE PLAN")
vline(s, 6.55, 1.9, 4.6, LINE, 1.0)
lh = box(s, ML, 1.9, 5.3, 0.4)
para(lh, "PIPELINE", 11, font=MONO, color=ACCENT, bold=True, first=True, tracking=1.6)
line(s, ML, 2.32, 5.3, LINE, 1.0)
lf = box(s, ML, 2.5, 5.3, 3.8)
for i, runs in enumerate([
    [("CLR(microbiome), log1p+standardise(metabolome); scaler & PCA fit on ", {}), ("train only", {"bold": True}), (".", {})],
    [("Core: AGW-OT and CM-Harmonic-Aniso vs. the baseline ladder.", {})],
    [("Ablations: ProMA (info-geometry floor) and DLB (diffusion-kernel check) to test what anchoring buys.", {})],
    [("Baseline ladder: mean / KNN / RF / PLS→frozen-PCA.", {})]]):
    para(lf, [("— ", {"color": ACCENT})] + runs, 13, font=DISP, color=INK2, spacing=1.28, space_after=9, first=(i == 0))
rh = box(s, 6.9, 1.9, 5.4, 0.4)
para(rh, "EVALUATION", 11, font=MONO, color=ACCENT, bold=True, first=True, tracking=1.6)
line(s, 6.9, 2.32, 5.3, LINE, 1.0)
rf = box(s, 6.9, 2.5, 5.5, 3.8)
for i, runs in enumerate([
    [("Mantel r vs. PCA-2PC ground truth (reproduce the grader exactly).", {})],
    [("Sweep missing-% = 10/30/50, ", {}), ("both directions", {"bold": True}), (".", {})],
    [("Subject-grouped CV", {"color": ACCENT, "bold": True}), (" — distance pairs are non-iid; no leakage.", {})],
    [("Report performance and runtime vs. the baseline.", {})]]):
    para(rf, [("— ", {"color": ACCENT})] + runs, 13, font=DISP, color=INK2, spacing=1.28, space_after=9, first=(i == 0))

# ---- 15 appendix ----
s = new(); shell(s, "Appendix", 15, num_pos="none")
eyebrow(s, ML, 1.05, "APPENDIX · SUPPORTING EXPLORATION")
pic(s, "06_between_group_distance.png", ML, 1.6, h=2.45)
pic(s, "08_metabolome_trends.png", 7.0, 1.6, h=2.45)
pic(s, "03_alpha_diversity.png", ML, 4.35, h=2.2)
pic(s, "13_gender_by_group.png", 7.0, 4.35, h=2.2)
for cx, cy, cap in [(ML, 4.12, "avg microbiome distance between groups"),
                    (7.0, 4.12, "metabolites elevated / reduced by group"),
                    (ML, 6.62, "Shannon diversity by group"),
                    (7.0, 6.62, "gender composition by group")]:
    cb = box(s, cx, cy, 5.6, 0.3)
    para(cb, cap.upper(), 8, font=MONO, color=MUT, first=True, tracking=0.6)

# ---- 16 closing ----
s = new(); slide_bg(s)
nt = box(s, 8.7, -0.4, 5.6, 3.5)
para(nt, "16", 250, font=DISP, color=ROSE_GHOST, bold=True, first=True, align=PP_ALIGN.RIGHT)
rect(s, ML, 0.46, 0.10, 0.10, ACCENT)
ht = box(s, ML + 0.22, 0.4, 7.0, 0.3)
para(ht, "METACARDIS · MILESTONE 1", 9.5, font=MONO, color=MUT, first=True, tracking=1.6)
line(s, ML, 0.82, W.inches - ML - MR, LINE, 1.0)
eyebrow(s, ML, 2.2, "KEY CHALLENGES → NEXT STEPS")
tf = box(s, ML, 2.7, 11.0, 1.8)
para(tf, [("Beat mean imputation by exploiting the ", {}),
          ("predictable cross-modal subset", {"color": ACCENT, "bold": True}),
          (" + the metabolome-only direction.", {})], 30, font=DISP, color=INK, first=True, spacing=1.05)
rect(s, ML, 4.9, 1.15, 0.05, ACCENT)
bf = box(s, ML, 5.2, 11.0, 0.8)
para(bf, "Next: implement AGW-OT + CM-Harmonic-Aniso · per-axis Mantel-vs-missing% sweep vs. the baseline ladder.",
     12, font=MONO, color=MUT, first=True, tracking=0.4)

prs.save(str(ROOT / "deliverables" / "milestone1.pptx"))
print("saved deliverables/milestone1.pptx with", len(prs.slides._sldIdLst), "slides")
